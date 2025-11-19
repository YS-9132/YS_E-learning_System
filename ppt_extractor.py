"""
PPTXファイルから問題を自動抽出するモジュール
スライドの構造に基づいて問題文・選択肢・正解を抽出
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from typing import List, Dict, Any


class PPTExtractor:
    """PPTXファイルから問題を抽出するクラス"""
    
    def __init__(self, pptx_path: str):
        """
        Args:
            pptx_path: PPTXファイルのパス
        """
        self.pptx_path = pptx_path
        self.presentation = None
        self.questions = []
        
    def load_presentation(self) -> bool:
        """PPTXファイルを読み込む"""
        try:
            self.presentation = Presentation(self.pptx_path)
            return True
        except Exception as e:
            print(f"エラー: PPTXファイルの読み込みに失敗しました: {e}")
            return False
    
    def extract_text_from_slide(self, slide) -> str:
        """スライドからすべてのテキストを抽出"""
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text.strip())
        return "\n".join(text)
    
    def extract_questions(self) -> List[Dict[str, Any]]:
        """
        PPTXから問題を抽出
        構造：スライド2以降が問題（スライド1はタイトル）
        """
        if not self.presentation:
            return []
        
        self.questions = []
        
        # スライド2以降を処理（スライド1はタイトルスライド）
        for slide_idx, slide in enumerate(self.presentation.slides[1:], start=1):
            text = self.extract_text_from_slide(slide)
            if not text.strip():
                continue
            
            question_data = self._parse_question_text(text, slide_idx)
            if question_data:
                self.questions.append(question_data)
        
        return self.questions
    
    def _parse_question_text(self, text: str, question_num: int) -> Dict[str, Any]:
        """
        テキストから問題データを抽出
        形式: 「問題N 問題文」+ 「A.～E. 選択肢」で○✕記載
        """
        lines = text.split('\n')
        
        question_text = ""
        choices = []
        correct_answers = []
        
        # 問題文と選択肢を分離
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 問題文の検出（「問題」で始まるか、数字で始まる）
            if line.startswith('問題') or (len(line) > 0 and line[0].isdigit() and '.' in line):
                # 問題番号を削除して問題文を取得
                question_text = re.sub(r'^問題\d+\s*', '', line)
                question_text = re.sub(r'^\d+\.\s*', '', question_text)
            
            # 選択肢の検出（A. B. C. D. E.）
            match = re.match(r'^([A-E])\.\s*(.+?)(\s*[○×✓✕])?$', line)
            if match:
                choice_letter = match.group(1)
                choice_text = match.group(2).strip()
                correct_mark = match.group(3).strip() if match.group(3) else ""
                
                # ○またはTrueで正解判定
                is_correct = "○" in correct_mark or "✓" in correct_mark
                
                choices.append({
                    "letter": choice_letter,
                    "text": choice_text,
                    "is_correct": is_correct
                })
                
                if is_correct:
                    correct_answers.append(choice_letter)
        
        if not question_text or not choices:
            return None
        
        return {
            "id": question_num,
            "question": question_text,
            "choices": choices,
            "correct_answers": correct_answers,
            "multiple_choice": len(correct_answers) > 1
        }
    
    def save_questions_to_json(self, output_path: str) -> bool:
        """問題をJSONファイルに保存"""
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(self.questions, f, ensure_ascii=False, indent=2)
            print(f"✅ 問題を保存しました: {output_path}")
            print(f"   {len(self.questions)}問を抽出しました")
            return True
        except Exception as e:
            print(f"エラー: JSONファイルの保存に失敗しました: {e}")
            return False
    
    def display_questions(self):
        """抽出された問題を表示（デバッグ用）"""
        for q in self.questions:
            print(f"\n{'='*60}")
            print(f"問題 {q['id']}: {q['question']}")
            print(f"複数選択: {q['multiple_choice']}")
            for choice in q['choices']:
                mark = "✅" if choice['is_correct'] else "  "
                print(f"  {mark} {choice['letter']}. {choice['text']}")
            print(f"正解: {', '.join(q['correct_answers'])}")


def main():
    """メイン処理（スタンドアロン実行用）"""
    import sys
    
    if len(sys.argv) < 2:
        print("使用方法: python ppt_extractor.py <pptxファイルパス> [出力パス]")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "questions.json"
    
    extractor = PPTExtractor(pptx_file)
    
    if not extractor.load_presentation():
        sys.exit(1)
    
    questions = extractor.extract_questions()
    
    if not questions:
        print("警告: 問題が抽出されませんでした")
        sys.exit(1)
    
    extractor.display_questions()
    extractor.save_questions_to_json(output_file)


if __name__ == "__main__":
    main()
